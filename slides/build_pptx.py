"""
Build slides/nexus-deck.pptx from scratch using python-pptx.

Why not Marp: marp-cli has a heavy network install (chromium) and was
failing on the build host. python-pptx is self-contained, ~2 MB, no
network needed. The trade-off is more imperative slide-construction
code here, but the design tokens stay identical to the Marp source
(slides/nexus-deck.md).

Run: python slides/build_pptx.py
Output: slides/nexus-deck.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------- Brand tokens
BG          = RGBColor(0x0E, 0x0E, 0x10)   # near-black background
PANEL       = RGBColor(0x1A, 0x1A, 0x1F)   # slightly lighter cards
PANEL_BORDER= RGBColor(0x2A, 0x2A, 0x2F)
FG          = RGBColor(0xF8, 0xF8, 0xF4)   # body text
FG_MUTED    = RGBColor(0xA0, 0xA0, 0xA0)
FG_FAINT    = RGBColor(0x55, 0x55, 0x55)
ACCENT_GOLD = RGBColor(0xFC, 0xD3, 0x4D)   # primary highlight
ACCENT_ORG  = RGBColor(0xE8, 0x6F, 0x2C)   # secondary highlight
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)   # tertiary

# Slide size — 16:9, 13.33in × 7.5in
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_deck() -> Presentation:
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def blank_slide(prs: Presentation):
    """Add a slide with no built-in layout content. Background = BG."""
    layout = prs.slide_layouts[6]  # 6 = blank
    slide = prs.slides.add_slide(layout)
    # Paint background near-black
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def add_text(slide, left, top, width, height, text, *,
             font_name="Calibri", size=18, bold=False,
             color=FG, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    """Drop a single-run text box at an absolute position."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font_name
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    return tb


def add_multiline(slide, left, top, width, height, lines, *,
                  font_name="Calibri", size=18,
                  color=FG, line_spacing=1.4,
                  align=PP_ALIGN.LEFT):
    """lines: list of dicts {text, bold?, color?, size?}"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line.get("text", "")
        f = run.font
        f.name = line.get("font", font_name)
        f.size = Pt(line.get("size", size))
        f.bold = line.get("bold", False)
        f.color.rgb = line.get("color", color)
    return tb


def add_panel(slide, left, top, width, height, *,
              fill=PANEL, border=PANEL_BORDER, accent_top=None):
    """Card-like rectangle with optional gold accent stripe on top."""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.04  # subtle rounding
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = border
    rect.line.width = Pt(0.5)
    rect.shadow.inherit = False
    if accent_top:
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent_top
        stripe.line.fill.background()
    return rect


def add_footer(slide, num: int, total: int, label: str = ""):
    text = f"NEXUS  ·  {label}  ·  {num} / {total}" if label else f"NEXUS  ·  {num} / {total}"
    add_text(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
             text, size=10, color=FG_FAINT)


# ============================================================
# Build each slide
# ============================================================

def slide_1_title(prs):
    s = blank_slide(prs)
    # Top brand strip
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Pt(6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_ORG; bar.line.fill.background()

    # NEXUS title — huge, centered
    add_text(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.6),
             "NEXUS", font_name="Calibri", size=110, bold=True,
             color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    # Tagline
    add_text(s, Inches(0.5), Inches(3.05), Inches(12.3), Inches(0.6),
             "The Living Agent Economy on Kite Chain",
             size=28, color=FG, align=PP_ALIGN.CENTER)
    # Hackathon meta
    add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.4),
             "Kite AI Global Hackathon 2026  ·  Novel Track",
             size=18, bold=True, color=ACCENT_ORG, align=PP_ALIGN.CENTER)

    # Live URL + GitHub + chain
    add_multiline(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.6),
        [
            {"text": "Live demo  →  https://44-215-246-131.nip.io", "size": 18, "color": FG, "bold": True},
            {"text": "GitHub  →  github.com/tarun-khatri/nexus-kiteAI-hackathon", "size": 16, "color": FG_MUTED},
            {"text": "Network  →  Kite Aero Testnet (Chain ID 2368)", "size": 16, "color": FG_MUTED},
        ], align=PP_ALIGN.CENTER, line_spacing=1.6)

    add_footer(s, 1, 10, "Title")


def slide_2_problem(prs):
    s = blank_slide(prs)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "THE PROBLEM", size=14, bold=True, color=ACCENT_ORG)
    add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
             "Agents work alone today.", size=46, bold=True, color=ACCENT_GOLD)

    add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.7),
             "Billions of dollars have gone into AI agent startups.  LangChain. AutoGPT. CrewAI. Adept. ChatGPT Plugins.",
             size=18, color=FG)

    add_text(s, Inches(0.7), Inches(3.1), Inches(12), Inches(0.5),
             "Yet not a single AI agent in production has ever paid another agent — on-chain, for a real service.",
             size=18, bold=True, color=FG)

    add_text(s, Inches(0.7), Inches(4.05), Inches(12), Inches(0.4),
             "They have:", size=18, color=FG_MUTED)

    bullets = [
        ("No shared", "identity"),
        ("No shared", "payment rail"),
        ("No shared", "rules"),
        ("No shared", "reputation"),
    ]
    y = Inches(4.55)
    for prefix, key in bullets:
        add_text(s, Inches(1.2), y, Inches(11), Inches(0.4),
                 f"·   {prefix}  ", size=18, color=FG_MUTED)
        add_text(s, Inches(3.2), y, Inches(9), Inches(0.4),
                 key, size=20, bold=True, color=ACCENT_GOLD)
        y += Inches(0.45)

    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.45),
             "So they can't form economies. They sit as isolated tools, wired by hand into one company's server.",
             size=16, color=FG_MUTED)

    add_footer(s, 2, 10, "The problem")


def slide_3_why_now(prs):
    s = blank_slide(prs)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "WHY NOW", size=14, bold=True, color=ACCENT_ORG)
    add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
             "Kite gives agents the missing primitives.", size=40, bold=True, color=ACCENT_GOLD)

    add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.5),
             "The first chain purpose-built for autonomous AI agents:",
             size=18, color=FG)

    # 4 primitive cards in a 2×2 grid
    primitives = [
        ("Agent Passport", "A real on-chain identity\n(queryable as a W3C-style DID)"),
        ("x402", "The native standard for paying\nmachines per call"),
        ("Verified Intent", "Programmable spending rules\na human signs"),
        ("On-chain Reputation", "A credit history\nbad actors can't fake"),
    ]
    card_w = Inches(5.8)
    card_h = Inches(1.4)
    start_x = Inches(0.7)
    start_y = Inches(3.05)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)

    for i, (title, desc) in enumerate(primitives):
        r = i // 2
        c = i % 2
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        add_panel(s, x, y, card_w, card_h, accent_top=ACCENT_ORG)
        add_text(s, x + Inches(0.2), y + Inches(0.15), card_w - Inches(0.4), Inches(0.4),
                 title, size=18, bold=True, color=ACCENT_GOLD)
        add_text(s, x + Inches(0.2), y + Inches(0.62), card_w - Inches(0.4), Inches(0.8),
                 desc, size=14, color=FG_MUTED)

    add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.5),
             "NEXUS is the first live system that uses all four — wired together as one running economy.",
             size=16, bold=True, color=FG)

    add_footer(s, 3, 10, "Why now")


def slide_4_solution(prs):
    s = blank_slide(prs)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "THE SOLUTION", size=14, bold=True, color=ACCENT_ORG)
    add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
             "NEXUS — an open economy of AI agents.", size=40, bold=True, color=ACCENT_GOLD)

    add_text(s, Inches(0.7), Inches(2.15), Inches(12), Inches(0.5),
             "11 AI agents.  23 capabilities.  One live marketplace on Kite.",
             size=20, bold=True, color=FG)

    bullets = [
        ("Agents", " discover each other through a public capability registry."),
        ("LLM router", " dynamically picks the best agents per query — no hardcoded pipelines."),
        ("Every job", " runs under an ECDSA-signed mandate with a 7-check circuit breaker."),
        ("Agents pay each other", " on-chain via x402, settled on Kite Aero Testnet."),
        ("Every output", " is audit-trailed on-chain with a SHA-256 traceability hash."),
        ("Market Pulse", " runs the whole economy autonomously — no human in the loop."),
    ]
    y = Inches(3.1)
    for bold_part, rest in bullets:
        # bullet dot
        add_text(s, Inches(0.95), y, Inches(0.3), Inches(0.4),
                 "•", size=18, bold=True, color=ACCENT_ORG)
        # bold + rest
        tb = s.shapes.add_textbox(Inches(1.3), y, Inches(11), Inches(0.45))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(0); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = bold_part
        r1.font.name = "Calibri"; r1.font.size = Pt(16); r1.font.bold = True
        r1.font.color.rgb = ACCENT_GOLD
        r2 = p.add_run(); r2.text = rest
        r2.font.name = "Calibri"; r2.font.size = Pt(16); r2.font.color.rgb = FG
        y += Inches(0.5)

    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             "Anyone can register a new agent in a single API call. The marketplace grows itself.",
             size=14, color=FG_MUTED)

    add_footer(s, 4, 10, "Solution")


def slide_5_traction(prs):
    s = blank_slide(prs)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "LIVE IN PRODUCTION", size=14, bold=True, color=ACCENT_ORG)
    add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
             "Real numbers. Not slideware.", size=40, bold=True, color=ACCENT_GOLD)

    # 6 stat cards in a 3×2 grid
    stats = [
        ("1,293", "On-chain transactions"),
        ("$0.1487", "USDC settled"),
        ("906", "Autonomous runs"),
        ("11", "Agents in the economy"),
        ("23", "Capabilities offered"),
        ("76", "Tests passing"),
    ]
    card_w = Inches(3.85)
    card_h = Inches(1.6)
    start_x = Inches(0.7)
    start_y = Inches(2.45)
    gap_x = Inches(0.2)
    gap_y = Inches(0.25)

    for i, (num, label) in enumerate(stats):
        r = i // 3
        c = i % 3
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        add_panel(s, x, y, card_w, card_h, accent_top=ACCENT_ORG)
        add_text(s, x + Inches(0.2), y + Inches(0.3), card_w - Inches(0.4), Inches(0.8),
                 num, size=42, bold=True, color=ACCENT_GOLD)
        add_text(s, x + Inches(0.2), y + Inches(1.05), card_w - Inches(0.4), Inches(0.4),
                 label.upper(), size=11, bold=True, color=FG_MUTED)

    add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.45),
             "Live URL  →  https://44-215-246-131.nip.io",
             size=18, bold=True, color=FG)
    add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.45),
             "Autonomous feed  →  https://44-215-246-131.nip.io/pulse",
             size=16, color=ACCENT_GOLD)
    add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
             "Every transaction above corresponds to a verifiable Kitescan tx hash.",
             size=12, color=FG_MUTED)

    add_footer(s, 5, 10, "Traction")


def slide_6_architecture(prs):
    s = blank_slide(prs)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "HOW IT WORKS", size=14, bold=True, color=ACCENT_ORG)
    add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
             "One pipeline. Eight on-chain steps.", size=40, bold=True, color=ACCENT_GOLD)

    steps = [
        ("1", "Query in", "User query OR Market Pulse autonomous tick"),
        ("2", "Discover", "LLM router reads live capability registry"),
        ("3", "Pick", "Best provider per capability (rep ↑ price ↓)"),
        ("4", "Sign", "Mandate Manager signs spending mandate (ECDSA)"),
        ("5", "Gate", "Circuit Breaker — 7 checks before every payment"),
        ("6", "Pay", "x402 micropayment fires on Kite Aero Testnet"),
        ("7", "Execute", "Orchestrator invokes the agent over its endpoint"),
        ("8", "Audit", "AuditAgent verifies → reputation + audit trail on-chain"),
    ]
    card_w = Inches(2.95)
    card_h = Inches(1.4)
    start_x = Inches(0.7)
    start_y = Inches(2.35)
    gap_x = Inches(0.13)
    gap_y = Inches(0.2)

    for i, (num, title, desc) in enumerate(steps):
        r = i // 4
        c = i % 4
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        add_panel(s, x, y, card_w, card_h, accent_top=ACCENT_GOLD)
        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.18),
                                  Inches(0.4), Inches(0.4))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT_ORG
        circ.line.fill.background()
        n_tb = circ.text_frame
        n_tb.margin_left = Pt(0); n_tb.margin_right = Pt(0)
        n_tb.margin_top = Pt(0); n_tb.margin_bottom = Pt(0)
        n_p = n_tb.paragraphs[0]; n_p.alignment = PP_ALIGN.CENTER
        n_r = n_p.add_run(); n_r.text = num
        n_r.font.name = "Calibri"; n_r.font.size = Pt(14); n_r.font.bold = True
        n_r.font.color.rgb = BG
        # title
        add_text(s, x + Inches(0.75), y + Inches(0.2), card_w - Inches(0.85), Inches(0.4),
                 title, size=15, bold=True, color=ACCENT_GOLD)
        # desc
        add_text(s, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), Inches(0.7),
                 desc, size=11, color=FG_MUTED)

    add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.5),
             "Steps 4, 5, 6, 8 are on-chain.  Every step is observable in real time on /pulse.",
             size=14, color=FG)

    add_footer(s, 6, 10, "Architecture")


def slide_7_autonomy(prs):
    s = blank_slide(prs)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "AGENT AUTONOMY", size=14, bold=True, color=ACCENT_ORG)
    add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
             "Market Pulse runs while we're talking.", size=38, bold=True, color=ACCENT_GOLD)

    add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.5),
             "Every hour, the backend wakes itself up — no human involvement.",
             size=18, color=FG)

    steps = [
        ("1", "Pull live market signals", "BTC/ETH/SOL 24h delta · Fear & Greed Index · CoinGecko trending"),
        ("2", "LLM-generate a fresh query", "Based on current conditions — never from a hardcoded list"),
        ("3", "Drive the full pipeline", "Mandate · x402 payments · audit trail · reputation update"),
        ("4", "Persist to /pulse", "Drillable to ECDSA signature, per-payment from→to, Kitescan tx hash"),
    ]
    y = Inches(3.1)
    for num, title, desc in steps:
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y, Inches(0.42), Inches(0.42))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT_ORG; circ.line.fill.background()
        n_tb = circ.text_frame
        n_tb.margin_left = Pt(0); n_tb.margin_right = Pt(0)
        n_tb.margin_top = Pt(0); n_tb.margin_bottom = Pt(0)
        n_p = n_tb.paragraphs[0]; n_p.alignment = PP_ALIGN.CENTER
        n_r = n_p.add_run(); n_r.text = num
        n_r.font.name = "Calibri"; n_r.font.size = Pt(14); n_r.font.bold = True
        n_r.font.color.rgb = BG

        add_text(s, Inches(1.5), y - Inches(0.02), Inches(11), Inches(0.4),
                 title, size=18, bold=True, color=ACCENT_GOLD)
        add_text(s, Inches(1.5), y + Inches(0.4), Inches(11), Inches(0.4),
                 desc, size=13, color=FG_MUTED)
        y += Inches(0.85)

    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
             "906 autonomous runs to date.  Bookmark /pulse and watch run #907 land on its own.",
             size=18, bold=True, color=ACCENT_GOLD)

    add_footer(s, 7, 10, "Autonomy")


def slide_8_devex(prs):
    s = blank_slide(prs)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "DEVELOPER EXPERIENCE", size=14, bold=True, color=ACCENT_ORG)
    add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
             "A new agent joins in one API call.", size=38, bold=True, color=ACCENT_GOLD)

    # Code card
    code_panel = add_panel(s, Inches(0.7), Inches(2.2), Inches(7.8), Inches(3.6),
                           fill=RGBColor(0x14, 0x14, 0x18), border=PANEL_BORDER,
                           accent_top=ACCENT_GOLD)
    code_lines = [
        ("POST /api/marketplace/register", ACCENT_GOLD, True),
        ("{", FG_MUTED, False),
        ("  \"name\": \"MyCustom-Agent-v1\",", FG, False),
        ("  \"description\": \"Tracks Solana memecoins...\",", FG, False),
        ("  \"capabilities\": [\"memecoin_discovery\"],", FG, False),
        ("  \"callback_url\": \"https://my-agent.example.com/invoke\",", FG, False),
        ("  \"price_per_query\": 0.0001,", FG, False),
        ("  \"keywords\": [\"memecoin\", \"solana\", \"trending\"],", FG, False),
        ("  \"example_queries\": [\"new Solana memes today\"]", FG, False),
        ("}", FG_MUTED, False),
    ]
    tb = s.shapes.add_textbox(Inches(0.95), Inches(2.4), Inches(7.4), Inches(3.4))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_top = Pt(0)
    for i, (line, color, bold) in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r = p.add_run(); r.text = line
        r.font.name = "Consolas"; r.font.size = Pt(13)
        r.font.color.rgb = color; r.font.bold = bold

    # Right column: benefits
    benefits = [
        ("Live in the economy", "the moment the call returns. Registry rebuilds."),
        ("Any language.", "Python, Node, Go, Rust — just expose an HTTP endpoint."),
        ("No code review.", "No team merge. No proxy permission to negotiate."),
        ("Real settlement", "from minute zero — every job is a real Kite tx."),
    ]
    y = Inches(2.4)
    for bold_part, rest in benefits:
        add_text(s, Inches(9.0), y, Inches(0.3), Inches(0.4),
                 "→", size=18, bold=True, color=ACCENT_ORG)
        tb = s.shapes.add_textbox(Inches(9.35), y, Inches(3.7), Inches(0.85))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(0); tf.margin_top = Pt(0)
        p = tf.paragraphs[0]; p.line_spacing = 1.35
        r1 = p.add_run(); r1.text = bold_part
        r1.font.name = "Calibri"; r1.font.size = Pt(14); r1.font.bold = True
        r1.font.color.rgb = ACCENT_GOLD
        r2 = p.add_run(); r2.text = " " + rest
        r2.font.name = "Calibri"; r2.font.size = Pt(13); r2.font.color.rgb = FG
        y += Inches(0.78)

    add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
             "The capability registry rebuilds, the LLM router picks the new agent up, x402 payments start landing.",
             size=14, color=FG_MUTED)

    add_footer(s, 8, 10, "DevEx")


def slide_9_diff(prs):
    s = blank_slide(prs)
    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "DIFFERENTIATION", size=14, bold=True, color=ACCENT_ORG)
    add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
             "Four things no other agent framework has.", size=36, bold=True, color=ACCENT_GOLD)

    # Comparison table
    headers = ["", "LangChain / CrewAI", "ChatGPT Plugins", "Coinbase AgentKit", "NEXUS"]
    rows = [
        ("Open marketplace", "✗", "✗", "✗", "✓"),
        ("Native machine payments (x402)", "✗", "✗", "partial", "✓ on-chain"),
        ("Verified Intent (ECDSA mandates)", "✗", "✗", "✗", "✓ 7-check"),
        ("On-chain reputation", "✗", "✗", "✗", "✓"),
        ("Autonomous trigger, no human", "✗", "✗", "✗", "✓"),
        ("Pure-LLM dynamic routing", "partial", "✗", "✗", "✓"),
    ]

    col_widths = [Inches(3.8), Inches(2.05), Inches(1.85), Inches(2.1), Inches(2.05)]
    row_h = Inches(0.42)
    start_x = Inches(0.7)
    start_y = Inches(2.15)

    # header row
    x = start_x
    for i, h in enumerate(headers):
        bg = PANEL if i < 4 else RGBColor(0x2A, 0x1F, 0x12)
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, col_widths[i], row_h)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        cell.line.color.rgb = PANEL_BORDER; cell.line.width = Pt(0.5)
        cell.shadow.inherit = False
        # header text
        if h:
            cell_tf = cell.text_frame
            cell_tf.margin_left = Pt(10); cell_tf.margin_right = Pt(10)
            cell_tf.margin_top = Pt(0); cell_tf.margin_bottom = Pt(0)
            cell_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = cell_tf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
            cr = cp.add_run(); cr.text = h
            cr.font.name = "Calibri"; cr.font.size = Pt(12); cr.font.bold = True
            cr.font.color.rgb = ACCENT_GOLD if i == 4 else FG
        x += col_widths[i]

    # data rows
    for r_idx, row in enumerate(rows):
        y = start_y + (r_idx + 1) * row_h
        x = start_x
        for c_idx, val in enumerate(row):
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_widths[c_idx], row_h)
            row_bg = BG if r_idx % 2 == 0 else RGBColor(0x12, 0x12, 0x16)
            if c_idx == 4:
                row_bg = RGBColor(0x1B, 0x14, 0x0A)  # subtle gold tint
            cell.fill.solid(); cell.fill.fore_color.rgb = row_bg
            cell.line.color.rgb = PANEL_BORDER; cell.line.width = Pt(0.3)
            cell.shadow.inherit = False
            cell_tf = cell.text_frame
            cell_tf.margin_left = Pt(10); cell_tf.margin_right = Pt(10)
            cell_tf.margin_top = Pt(0); cell_tf.margin_bottom = Pt(0)
            cell_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cp = cell_tf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            cr = cp.add_run(); cr.text = val
            cr.font.name = "Calibri"; cr.font.size = Pt(12)
            if c_idx == 0:
                cr.font.color.rgb = FG
            elif val.startswith("✓"):
                cr.font.color.rgb = ACCENT_GOLD; cr.font.bold = True
            elif val == "✗":
                cr.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
            else:  # "partial"
                cr.font.color.rgb = ACCENT_ORG
            x += col_widths[c_idx]

    add_text(s, Inches(0.7), Inches(5.85), Inches(12), Inches(0.4),
             "Built on free-tier infrastructure  ·  Groq + Gemini fallback + Kite testnet + AWS t3.small",
             size=14, color=FG)
    add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.5),
             "Total cost to operate in production:  $0 / month",
             size=20, bold=True, color=ACCENT_GOLD)

    add_footer(s, 9, 10, "Differentiation")


def slide_10_vision(prs):
    s = blank_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Pt(6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_ORG; bar.line.fill.background()

    add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4),
             "WHAT COMES NEXT", size=14, bold=True, color=ACCENT_ORG)
    add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
             "Roadmap.", size=44, bold=True, color=ACCENT_GOLD)

    milestones = [
        ("Now",   "Live on Kite Aero Testnet  ·  open source  ·  906 autonomous runs"),
        ("+30d",  "Subscription mandates  ·  recurring x402 on signed schedules"),
        ("+60d",  "NEXUS Reputation Oracle  ·  any dApp on Kite queries agent rep on-chain"),
        ("+90d",  "Agent SDK  ·  one-line marketplace plug-in for any builder on Kite"),
        ("Mainnet", "Real USDC settlement when Kite mainnet ships"),
    ]
    y = Inches(2.15)
    for tag, desc in milestones:
        # tag pill
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y,
                                  Inches(1.5), Inches(0.4))
        pill.adjustments[0] = 0.4
        pill.fill.solid(); pill.fill.fore_color.rgb = ACCENT_ORG
        pill.line.fill.background()
        ptf = pill.text_frame
        ptf.margin_left = Pt(0); ptf.margin_right = Pt(0)
        ptf.margin_top = Pt(0); ptf.margin_bottom = Pt(0)
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run(); pr.text = tag
        pr.font.name = "Calibri"; pr.font.size = Pt(13); pr.font.bold = True
        pr.font.color.rgb = BG

        add_text(s, Inches(2.4), y + Inches(0.06), Inches(10.4), Inches(0.4),
                 desc, size=15, color=FG)
        y += Inches(0.55)

    # Tagline
    add_panel(s, Inches(0.7), Inches(5.15), Inches(12.0), Inches(1.2),
              fill=RGBColor(0x16, 0x14, 0x10), border=ACCENT_ORG)
    add_text(s, Inches(1.0), Inches(5.35), Inches(11.5), Inches(0.5),
             "Stripe gave the internet a payment rail.  Yelp gave it reputation.  Upwork gave it a marketplace.",
             size=15, color=FG_MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(5.78), Inches(11.5), Inches(0.5),
             "NEXUS gives the agent economy all three — on Kite.",
             size=20, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

    # Final URLs
    add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
             "Live:  https://44-215-246-131.nip.io     ·     GitHub:  github.com/tarun-khatri/nexus-kiteAI-hackathon",
             size=13, color=FG, align=PP_ALIGN.CENTER)

    add_footer(s, 10, 10, "Vision")


# ============================================================
# Build & save
# ============================================================

def build(out_path: Path) -> None:
    prs = new_deck()
    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_why_now(prs)
    slide_4_solution(prs)
    slide_5_traction(prs)
    slide_6_architecture(prs)
    slide_7_autonomy(prs)
    slide_8_devex(prs)
    slide_9_diff(prs)
    slide_10_vision(prs)
    prs.save(str(out_path))
    print(f"[OK] Wrote {out_path}  ({out_path.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    here = Path(__file__).parent
    build(here / "nexus-deck.pptx")
